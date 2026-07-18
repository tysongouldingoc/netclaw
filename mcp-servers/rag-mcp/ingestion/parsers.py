"""Two-tier document parsing for rag-mcp (FR-002, FR-008a).

Native tier: PDF (pymupdf), Markdown, HTML (BeautifulSoup), TXT.
Office tier: DOCX/XLSX/PPTX/VSDX via Python parsers; legacy DOC/XLS/PPT/VSD
via `soffice --headless --convert-to pdf` fallback (graceful failure with an
actionable message when LibreOffice is absent).

Output is a normalized ParsedDocument: a list of Section objects
(heading path + body blocks), page count where the format provides pages,
and the SHA-256 content hash used for dedupe (FR-006/FR-007).

Blocks are (kind, text, page) tuples where kind is 'text' | 'atomic'
('atomic' = fenced code/CLI/config block or table — never split, FR-011).
"""

import hashlib
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Tuple

NATIVE_EXTENSIONS = {".pdf", ".md", ".markdown", ".html", ".htm", ".txt"}
MODERN_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".vsdx"}
LEGACY_OFFICE_EXTENSIONS = {".doc", ".xls", ".ppt", ".vsd"}
SUPPORTED_EXTENSIONS = NATIVE_EXTENSIONS | MODERN_OFFICE_EXTENSIONS | LEGACY_OFFICE_EXTENSIONS


class IngestError(Exception):
    def __init__(self, code: str, message: str):
        self.code = code
        self.message = message
        super().__init__(f"{code}: {message}")


@dataclass
class Section:
    """A heading-scoped run of content. heading_path excludes the doc title."""

    heading_path: List[str]
    blocks: List[Tuple[str, str, Optional[int]]] = field(default_factory=list)  # (kind, text, page)


@dataclass
class ParsedDocument:
    title: str
    sections: List[Section]
    page_count: Optional[int]
    content_hash: str


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(1 << 20), b""):
            h.update(block)
    return h.hexdigest()


def check_size_cap(path: Path, max_mb: int, max_pages: int, page_count: Optional[int] = None) -> None:
    size_mb = path.stat().st_size / (1024 * 1024)
    if size_mb > max_mb:
        raise IngestError(
            "SIZE_LIMIT_EXCEEDED",
            f"Document is {size_mb:.1f} MB — over the {max_mb} MB cap. "
            f"Raise RAG_MAX_DOC_MB in .env to override.",
        )
    if page_count is not None and page_count > max_pages:
        raise IngestError(
            "SIZE_LIMIT_EXCEEDED",
            f"Document has {page_count} pages — over the {max_pages}-page cap. "
            f"Raise RAG_MAX_DOC_PAGES in .env to override.",
        )


class _HeadingStack:
    """Level-aware heading stack: pushing a heading pops everything at the
    same or deeper level, so sibling headings replace rather than nest —
    regardless of level gaps or a consumed H1 title."""

    def __init__(self):
        self._stack: List[Tuple[int, str]] = []

    def push(self, level: int, text: str) -> List[str]:
        while self._stack and self._stack[-1][0] >= level:
            self._stack.pop()
        self._stack.append((level, text))
        return [t for _, t in self._stack]


def _require_text(sections: List[Section], source: str) -> None:
    if not any(blk[1].strip() for s in sections for blk in s.blocks):
        raise IngestError(
            "PARSE_FAILED",
            f"No extractable text in '{source}' — scanned/image-only documents "
            "are not supported (OCR is out of scope for v1).",
        )


# ---------------------------------------------------------------------
# PDF (pymupdf): TOC headings when present, font-size heuristic fallback
# ---------------------------------------------------------------------
def _parse_pdf(path: Path, max_pages: int) -> ParsedDocument:
    try:
        import fitz  # pymupdf
    except ImportError as exc:
        raise IngestError("PARSE_FAILED", f"pymupdf not installed: {exc}")

    try:
        doc = fitz.open(str(path))
    except Exception as exc:
        raise IngestError("PARSE_FAILED", f"Cannot open PDF '{path.name}': {exc}")

    page_count = doc.page_count
    if page_count > max_pages:
        doc.close()
        raise IngestError(
            "SIZE_LIMIT_EXCEEDED",
            f"Document has {page_count} pages — over the {max_pages}-page cap. "
            f"Raise RAG_MAX_DOC_PAGES in .env to override.",
        )

    title = (doc.metadata or {}).get("title") or path.stem

    # Heading map: page -> [(level, heading_text)] from TOC when available
    toc = doc.get_toc(simple=True) or []
    headings_by_page = {}
    for level, text, page_no in toc:
        headings_by_page.setdefault(page_no, []).append((level, text.strip()))

    sections: List[Section] = []
    current = Section(heading_path=[])
    sections.append(current)
    hstack = _HeadingStack()

    # Font-size heuristic for PDFs without a TOC: body size = most common size
    use_font_heuristic = not toc

    for page_index in range(page_count):
        page = doc[page_index]
        page_no = page_index + 1

        for level, htext in headings_by_page.get(page_no, []):
            current = Section(heading_path=hstack.push(level, htext))
            sections.append(current)

        if use_font_heuristic:
            blocks = page.get_text("dict")["blocks"]
            sizes = [
                round(span["size"])
                for b in blocks
                if b.get("type") == 0
                for line in b.get("lines", [])
                for span in line.get("spans", [])
            ]
            body_size = max(set(sizes), key=sizes.count) if sizes else 0
            for b in blocks:
                if b.get("type") != 0:
                    continue
                for line in b.get("lines", []):
                    text = "".join(s["text"] for s in line.get("spans", [])).strip()
                    if not text:
                        continue
                    line_size = max(round(s["size"]) for s in line["spans"])
                    if line_size > body_size + 1 and len(text) < 120:
                        current = Section(heading_path=hstack.push(1, text))
                        sections.append(current)
                    else:
                        current.blocks.append(("text", text, page_no))
        else:
            text = page.get_text("text").strip()
            if text:
                for para in re.split(r"\n\s*\n", text):
                    para = para.strip()
                    if para:
                        current.blocks.append(("text", para, page_no))

        # Tables are atomic blocks (FR-011)
        try:
            for table in page.find_tables().tables:
                rows = table.extract()
                rendered = "\n".join(
                    " | ".join("" if c is None else str(c) for c in row) for row in rows
                ).strip()
                if rendered:
                    current.blocks.append(("atomic", rendered, page_no))
        except Exception:
            pass

    doc.close()
    sections = [s for s in sections if s.blocks]
    _require_text(sections, path.name)
    return ParsedDocument(title=title, sections=sections, page_count=page_count, content_hash="")


# ---------------------------------------------------------------------
# Markdown: headings, fenced code blocks (atomic), pipe tables (atomic)
# ---------------------------------------------------------------------
_MD_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")
_MD_FENCE = re.compile(r"^(```|~~~)")
_MD_TABLE_ROW = re.compile(r"^\s*\|.*\|\s*$")


def _parse_markdown(text: str, title: str) -> ParsedDocument:
    sections: List[Section] = [Section(heading_path=[])]
    current = sections[0]
    hstack = _HeadingStack()
    seen_structural_heading = False
    doc_title = title

    lines = text.splitlines()
    i, n = 0, len(lines)
    para: List[str] = []

    def flush_para():
        nonlocal para
        joined = "\n".join(para).strip()
        if joined:
            current.blocks.append(("text", joined, None))
        para = []

    while i < n:
        line = lines[i]
        m = _MD_HEADING.match(line)
        if m:
            flush_para()
            level = len(m.group(1))
            htext = m.group(2).strip()
            if level == 1 and not any(s.blocks for s in sections) and not seen_structural_heading:
                doc_title = htext
                i += 1
                continue
            seen_structural_heading = True
            current = Section(heading_path=hstack.push(level, htext))
            sections.append(current)
            i += 1
            continue
        if _MD_FENCE.match(line):
            flush_para()
            fence = _MD_FENCE.match(line).group(1)
            block = [line]
            i += 1
            while i < n and not lines[i].startswith(fence):
                block.append(lines[i])
                i += 1
            if i < n:
                block.append(lines[i])
                i += 1
            current.blocks.append(("atomic", "\n".join(block), None))
            continue
        if _MD_TABLE_ROW.match(line):
            flush_para()
            block = []
            while i < n and _MD_TABLE_ROW.match(lines[i]):
                block.append(lines[i])
                i += 1
            current.blocks.append(("atomic", "\n".join(block), None))
            continue
        if not line.strip():
            flush_para()
            i += 1
            continue
        para.append(line)
        i += 1
    flush_para()

    sections = [s for s in sections if s.blocks]
    _require_text(sections, title)
    return ParsedDocument(title=doc_title, sections=sections, page_count=None, content_hash="")


# ---------------------------------------------------------------------
# HTML: heading hierarchy preserved, boilerplate stripped
# ---------------------------------------------------------------------
def parse_html_text(html: str, fallback_title: str) -> ParsedDocument:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise IngestError("PARSE_FAILED", f"beautifulsoup4 not installed: {exc}")

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    title = (soup.title.string.strip() if soup.title and soup.title.string else fallback_title)

    sections: List[Section] = [Section(heading_path=[])]
    current = sections[0]
    hstack = _HeadingStack()

    body = soup.body or soup
    for el in body.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "pre", "table", "li"]):
        if el.name.startswith("h"):
            level = int(el.name[1])
            htext = el.get_text(" ", strip=True)
            if not htext:
                continue
            current = Section(heading_path=hstack.push(level, htext))
            sections.append(current)
        elif el.name == "pre":
            text = el.get_text("\n", strip=True)
            if text:
                current.blocks.append(("atomic", text, None))
        elif el.name == "table":
            rows = [
                " | ".join(c.get_text(" ", strip=True) for c in tr.find_all(["td", "th"]))
                for tr in el.find_all("tr")
            ]
            rendered = "\n".join(r for r in rows if r.strip())
            if rendered:
                current.blocks.append(("atomic", rendered, None))
        else:
            # Skip text nested inside pre/table (already captured atomically)
            if el.find_parent(["pre", "table"]):
                continue
            text = el.get_text(" ", strip=True)
            if text:
                current.blocks.append(("text", text, None))

    sections = [s for s in sections if s.blocks]
    _require_text(sections, fallback_title)
    return ParsedDocument(title=title, sections=sections, page_count=None, content_hash="")


def _parse_txt(text: str, title: str) -> ParsedDocument:
    section = Section(heading_path=[])
    for para in re.split(r"\n\s*\n", text):
        para = para.strip()
        if para:
            section.blocks.append(("text", para, None))
    _require_text([section], title)
    return ParsedDocument(title=title, sections=[section], page_count=None, content_hash="")


# ---------------------------------------------------------------------
# Modern office tier
# ---------------------------------------------------------------------
def _parse_docx(path: Path) -> ParsedDocument:
    try:
        import docx
    except ImportError as exc:
        raise IngestError("PARSE_FAILED", f"python-docx not installed: {exc}")

    d = docx.Document(str(path))
    title = path.stem
    sections: List[Section] = [Section(heading_path=[])]
    current = sections[0]
    hstack = _HeadingStack()

    for para in d.paragraphs:
        style = (para.style.name or "") if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        m = re.match(r"Heading (\d+)", style)
        if style == "Title":
            title = text
        elif m:
            level = int(m.group(1))
            current = Section(heading_path=hstack.push(level, text))
            sections.append(current)
        else:
            current.blocks.append(("text", text, None))

    for table in d.tables:
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells) for row in table.rows
        ]
        rendered = "\n".join(r for r in rows if r.strip())
        if rendered:
            current.blocks.append(("atomic", rendered, None))

    sections = [s for s in sections if s.blocks]
    _require_text(sections, path.name)
    return ParsedDocument(title=title, sections=sections, page_count=None, content_hash="")


def _parse_xlsx(path: Path) -> ParsedDocument:
    try:
        import openpyxl
    except ImportError as exc:
        raise IngestError("PARSE_FAILED", f"openpyxl not installed: {exc}")

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sections: List[Section] = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            section = Section(heading_path=[sheet.title])
            section.blocks.append(("atomic", "\n".join(rows), None))
            sections.append(section)
    wb.close()
    _require_text(sections, path.name)
    return ParsedDocument(title=path.stem, sections=sections, page_count=len(sections), content_hash="")


def _parse_pptx(path: Path) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise IngestError("PARSE_FAILED", f"python-pptx not installed: {exc}")

    prs = Presentation(str(path))
    sections: List[Section] = []
    slide_count = 0
    for idx, slide in enumerate(prs.slides, start=1):
        slide_count = idx
        slide_title = None
        texts = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if slide_title is None and shape == getattr(slide.shapes, "title", None):
                slide_title = text
            else:
                texts.append(text)
        heading = slide_title or f"Slide {idx}"
        section = Section(heading_path=[heading])
        for t in texts:
            section.blocks.append(("text", t, idx))
        if slide_title and not texts:
            section.blocks.append(("text", slide_title, idx))
        if section.blocks:
            sections.append(section)
    _require_text(sections, path.name)
    return ParsedDocument(title=path.stem, sections=sections, page_count=slide_count, content_hash="")


def _parse_vsdx(path: Path) -> ParsedDocument:
    try:
        from vsdx import VisioFile
    except ImportError as exc:
        raise IngestError("PARSE_FAILED", f"vsdx not installed: {exc}")

    sections: List[Section] = []
    with VisioFile(str(path)) as vis:
        for idx, page in enumerate(vis.pages, start=1):
            texts = []
            for shape in page.all_shapes:
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
            if texts:
                section = Section(heading_path=[page.name or f"Page {idx}"])
                section.blocks.append(("text", "\n".join(texts), idx))
                sections.append(section)
    _require_text(sections, path.name)
    return ParsedDocument(title=path.stem, sections=sections, page_count=len(sections), content_hash="")


# ---------------------------------------------------------------------
# Legacy office tier: LibreOffice headless -> PDF -> PDF path
# ---------------------------------------------------------------------
def _convert_legacy(path: Path, max_pages: int) -> ParsedDocument:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise IngestError(
            "CONVERTER_UNAVAILABLE",
            f"Legacy format '{path.suffix}' needs LibreOffice for offline conversion. "
            "Install it with: sudo apt install libreoffice (or re-run the rag-mcp "
            "installer step and accept the LibreOffice option). Modern formats "
            "(PDF/MD/HTML/TXT/DOCX/XLSX/PPTX/VSDX) work without it.",
        )
    with tempfile.TemporaryDirectory() as tmp:
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(path)],
                check=True,
                capture_output=True,
                timeout=300,
            )
        except subprocess.SubprocessError as exc:
            raise IngestError("PARSE_FAILED", f"LibreOffice conversion failed for '{path.name}': {exc}")
        converted = Path(tmp) / (path.stem + ".pdf")
        if not converted.exists():
            raise IngestError("PARSE_FAILED", f"LibreOffice produced no output for '{path.name}'")
        parsed = _parse_pdf(converted, max_pages)
        return ParsedDocument(
            title=path.stem if parsed.title == converted.stem else parsed.title,
            sections=parsed.sections,
            page_count=parsed.page_count,
            content_hash="",
        )


# ---------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------
def parse_file(path: Path, max_mb: int = 100, max_pages: int = 1000) -> ParsedDocument:
    path = Path(path)
    if not path.exists():
        raise IngestError("NOT_FOUND", f"File not found: {path}")
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise IngestError(
            "UNSUPPORTED_FORMAT", f"'{ext}' is not supported. Supported formats: {supported}"
        )

    check_size_cap(path, max_mb, max_pages)
    content_hash = sha256_file(path)

    if ext == ".pdf":
        parsed = _parse_pdf(path, max_pages)
    elif ext in (".md", ".markdown"):
        parsed = _parse_markdown(path.read_text(encoding="utf-8", errors="replace"), path.stem)
    elif ext in (".html", ".htm"):
        parsed = parse_html_text(path.read_text(encoding="utf-8", errors="replace"), path.stem)
    elif ext == ".txt":
        parsed = _parse_txt(path.read_text(encoding="utf-8", errors="replace"), path.stem)
    elif ext == ".docx":
        parsed = _parse_docx(path)
    elif ext == ".xlsx":
        parsed = _parse_xlsx(path)
    elif ext == ".pptx":
        parsed = _parse_pptx(path)
    elif ext == ".vsdx":
        parsed = _parse_vsdx(path)
    else:  # legacy: .doc, .xls, .ppt, .vsd
        parsed = _convert_legacy(path, max_pages)

    parsed.content_hash = content_hash
    return parsed
